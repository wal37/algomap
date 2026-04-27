from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("/Users/wale/Desktop/geonex/geoweb/ALGOMAP_PITCH_DECK.pptx")


def rgb(value: str) -> RGBColor:
    value = value.replace("#", "")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


NIGHT = rgb("#0D1117")
SLATE = rgb("#161F2B")
PANEL = rgb("#1D2A39")
STEEL = rgb("#5D7387")
CYAN = rgb("#45D3E8")
GREEN = rgb("#6EE7A2")
AMBER = rgb("#F0B44A")
ROSE = rgb("#E7846D")
MIST = rgb("#D8E0EA")
FOG = rgb("#EDF2F7")
WHITE = rgb("#FFFFFF")
INK = rgb("#0E1720")
MUTED = rgb("#607286")
GRID = rgb("#243447")


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_round(slide, left, top, width, height, fill, line=None, radius=0.04):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.adjustments[0] = radius
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=20,
    color=INK,
    bold=False,
    font="Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_label(slide, left, top, text, fill=CYAN, color=INK):
    width = max(1.45, 0.102 * len(text) + 0.45)
    add_round(slide, left, top, Inches(width), Inches(0.38), fill, None, 0.14)
    add_text(slide, left + Inches(0.16), top + Inches(0.07), Inches(width - 0.18), Inches(0.2), text.upper(), size=10, color=color, bold=True)


def add_card(slide, left, top, width, height, title, body, fill, title_color=WHITE, body_color=MIST):
    add_round(slide, left, top, width, height, fill, None, 0.035)
    add_text(slide, left + Inches(0.22), top + Inches(0.22), width - Inches(0.38), Inches(0.34), title, size=18, color=title_color, bold=True)
    add_text(slide, left + Inches(0.22), top + Inches(0.72), width - Inches(0.4), height - Inches(0.92), body, size=12, color=body_color)


def add_logo(slide, left, top, size=0.78):
    scale = size / 0.78
    add_round(slide, left, top, Inches(0.78 * scale), Inches(0.78 * scale), SLATE, None, 0.16)
    add_rect(slide, left + Inches(0.12 * scale), top + Inches(0.14 * scale), Inches(0.12 * scale), Inches(0.52 * scale), CYAN)
    add_rect(slide, left + Inches(0.32 * scale), top + Inches(0.24 * scale), Inches(0.12 * scale), Inches(0.42 * scale), GREEN)
    add_rect(slide, left + Inches(0.52 * scale), top + Inches(0.08 * scale), Inches(0.12 * scale), Inches(0.58 * scale), AMBER)
    add_rect(slide, left + Inches(0.12 * scale), top + Inches(0.56 * scale), Inches(0.52 * scale), Inches(0.08 * scale), WHITE)


def add_page(slide, num, fill=NIGHT):
    add_rect(slide, 0, 0, Inches(13.333), Inches(7.5), fill)
    add_rect(slide, Inches(0.72), Inches(0.72), Inches(11.9), Inches(0.02), GRID if fill == NIGHT else STEEL)
    add_rect(slide, Inches(0.72), Inches(6.82), Inches(11.9), Inches(0.02), GRID if fill == NIGHT else STEEL)
    add_text(slide, Inches(12.0), Inches(0.24), Inches(0.42), Inches(0.18), f"{num:02d}", size=10, color=MIST if fill == NIGHT else INK, bold=True, align=PP_ALIGN.RIGHT)


def cover(slide, num):
    add_page(slide, num, fill=NIGHT)
    add_rect(slide, Inches(7.86), Inches(0.72), Inches(4.76), Inches(6.1), PANEL)
    add_rect(slide, Inches(8.2), Inches(1.1), Inches(4.08), Inches(1.02), CYAN)
    add_rect(slide, Inches(8.72), Inches(2.48), Inches(3.12), Inches(0.82), GREEN)
    add_rect(slide, Inches(8.2), Inches(3.72), Inches(4.08), Inches(2.56), AMBER)
    add_logo(slide, Inches(1.0), Inches(1.02), 0.84)
    add_label(slide, Inches(1.0), Inches(2.02), "GIS and logistics operating system")
    add_text(slide, Inches(1.0), Inches(2.62), Inches(5.8), Inches(0.82), "Algomap", size=36, color=WHITE, bold=True)
    add_text(slide, Inches(1.0), Inches(3.58), Inches(6.0), Inches(1.44), "A routing, mapping, and spatial operations platform for fleets, assets, field teams, and infrastructure decisions.", size=21, color=MIST)
    add_text(slide, Inches(1.0), Inches(6.16), Inches(2.8), Inches(0.22), "Pitch deck | April 2026", size=12, color=CYAN, bold=True)


def problem(slide, num):
    add_page(slide, num, fill=FOG)
    add_label(slide, Inches(1.0), Inches(1.02), "Problem", fill=ROSE, color=WHITE)
    add_text(slide, Inches(1.0), Inches(1.56), Inches(10.8), Inches(0.82), "Spatial operations still break down across disconnected maps, dispatch tools, records, and manual routing decisions.", size=29, color=INK, bold=True)
    add_rect(slide, Inches(1.0), Inches(2.68), Inches(11.1), Inches(0.06), INK)
    add_card(slide, Inches(1.0), Inches(3.16), Inches(3.46), Inches(2.18), "Fragmented map systems", "Routing, field data, infrastructure layers, and operations records often live in separate tools.", WHITE, title_color=INK, body_color=MUTED)
    add_card(slide, Inches(4.82), Inches(3.16), Inches(3.46), Inches(2.18), "Slow dispatch and planning", "Logistics teams lose time moving between schedules, maps, and incident records.", MIST, title_color=INK, body_color=MUTED)
    add_card(slide, Inches(8.64), Inches(3.16), Inches(3.46), Inches(2.18), "Weak historical visibility", "Location-linked decisions are hard to review, compare, and optimize over time.", WHITE, title_color=INK, body_color=MUTED)


def why_now(slide, num):
    add_page(slide, num, fill=WHITE)
    add_label(slide, Inches(1.0), Inches(1.02), "Why now", fill=GREEN, color=INK)
    add_text(slide, Inches(1.0), Inches(1.56), Inches(6.6), Inches(0.82), "Operational teams need one spatial layer that connects movement, assets, and decisions.", size=28, color=INK, bold=True)
    add_text(slide, Inches(1.0), Inches(2.5), Inches(6.3), Inches(0.94), "As logistics networks, field services, utilities, and infrastructure programs scale, the bottleneck shifts toward coordination across geography, not just more raw location data.", size=15, color=MUTED)
    add_rect(slide, Inches(7.72), Inches(1.0), Inches(4.2), Inches(5.3), PANEL)
    add_card(slide, Inches(8.02), Inches(1.34), Inches(3.6), Inches(1.0), "Shift 1", "Every operation is now location-linked.", CYAN, title_color=INK, body_color=INK)
    add_card(slide, Inches(8.02), Inches(2.84), Inches(3.6), Inches(1.0), "Shift 2", "Fleet, field, and infrastructure teams need faster routing decisions.", GREEN, title_color=INK, body_color=INK)
    add_card(slide, Inches(8.02), Inches(4.34), Inches(3.6), Inches(1.0), "Shift 3", "Historical spatial records are becoming a competitive asset.", AMBER, title_color=INK, body_color=INK)


def product(slide, num):
    add_page(slide, num, fill=NIGHT)
    add_label(slide, Inches(1.0), Inches(1.02), "Platform", fill=WHITE, color=INK)
    add_text(slide, Inches(1.0), Inches(1.56), Inches(6.8), Inches(0.82), "Algomap is a spatial operating system for routing, mapping, and logistics decisions.", size=28, color=WHITE, bold=True)
    add_text(slide, Inches(1.0), Inches(2.5), Inches(6.7), Inches(0.96), "The platform combines GIS layers, route planning, asset visibility, field records, and operational analytics into one working interface.", size=15, color=MIST)
    add_card(slide, Inches(1.0), Inches(4.14), Inches(2.58), Inches(1.74), "Mapping layer", "Base maps, boundaries, roads, assets, and live overlays.", CYAN, title_color=INK, body_color=INK)
    add_card(slide, Inches(3.86), Inches(4.14), Inches(2.58), Inches(1.74), "Routing layer", "Stops, optimization, ETA logic, and movement planning.", GREEN, title_color=INK, body_color=INK)
    add_card(slide, Inches(6.72), Inches(4.14), Inches(2.58), Inches(1.74), "Field layer", "Incident logs, inspections, photos, and crew activity.", AMBER, title_color=INK, body_color=INK)
    add_card(slide, Inches(9.58), Inches(4.14), Inches(2.58), Inches(1.74), "Analytics layer", "Performance, coverage, and spatial history.", WHITE, title_color=INK, body_color=MUTED)


def components(slide, num):
    add_page(slide, num, fill=FOG)
    add_label(slide, Inches(1.0), Inches(1.02), "Core components", fill=CYAN, color=INK)
    add_text(slide, Inches(1.0), Inches(1.56), Inches(6.8), Inches(0.82), "Built around the full spatial operations stack, not one isolated map screen.", size=28, color=INK, bold=True)
    add_card(slide, Inches(1.0), Inches(3.04), Inches(3.48), Inches(2.12), "GIS and layer management", "Administrative boundaries, infrastructure layers, service zones, asset points, and custom overlays.", WHITE, title_color=INK, body_color=MUTED)
    add_card(slide, Inches(4.84), Inches(3.04), Inches(3.48), Inches(2.12), "Routing and dispatch", "Route creation, optimization, scheduling, stop sequencing, and travel-time operations.", MIST, title_color=INK, body_color=MUTED)
    add_card(slide, Inches(8.68), Inches(3.04), Inches(3.48), Inches(2.12), "Field records and monitoring", "Inspections, incidents, work orders, spatial logs, and coverage tracking.", WHITE, title_color=INK, body_color=MUTED)


def logistics(slide, num):
    add_page(slide, num, fill=WHITE)
    add_label(slide, Inches(1.0), Inches(1.02), "Logistics", fill=AMBER, color=INK)
    add_text(slide, Inches(1.0), Inches(1.56), Inches(6.4), Inches(0.82), "Routing is only valuable when it stays connected to real operating conditions.", size=28, color=INK, bold=True)
    add_text(slide, Inches(1.0), Inches(2.46), Inches(6.2), Inches(0.96), "Algomap ties dispatch and route planning to geography, field updates, asset conditions, and historical movement patterns.", size=15, color=MUTED)
    add_rect(slide, Inches(7.78), Inches(0.96), Inches(4.18), Inches(5.6), SLATE)
    add_card(slide, Inches(8.08), Inches(1.3), Inches(3.58), Inches(1.0), "Route planning", "Plan stops, service zones, delivery paths, and time windows.", CYAN, title_color=INK, body_color=INK)
    add_card(slide, Inches(8.08), Inches(2.78), Inches(3.58), Inches(1.0), "Fleet visibility", "Track vehicles, crews, and coverage by area or assignment.", GREEN, title_color=INK, body_color=INK)
    add_card(slide, Inches(8.08), Inches(4.26), Inches(3.58), Inches(1.0), "Exception handling", "Respond to congestion, outages, delays, or field incidents from one map context.", AMBER, title_color=INK, body_color=INK)


def gis(slide, num):
    add_page(slide, num, fill=MIST)
    add_label(slide, Inches(1.0), Inches(1.02), "GIS layer", fill=SLATE, color=WHITE)
    add_text(slide, Inches(1.0), Inches(1.56), Inches(6.8), Inches(0.82), "The GIS layer turns maps into operational structure, not just visualization.", size=28, color=INK, bold=True)
    add_card(slide, Inches(1.0), Inches(3.1), Inches(2.68), Inches(1.94), "Asset mapping", "Roads, depots, service points, pipelines, towers, or delivery nodes.", WHITE, title_color=INK, body_color=MUTED)
    add_card(slide, Inches(3.94), Inches(3.1), Inches(2.68), Inches(1.94), "Zone management", "Coverage areas, districts, clusters, and assignment boundaries.", FOG, title_color=INK, body_color=MUTED)
    add_card(slide, Inches(6.88), Inches(3.1), Inches(2.68), Inches(1.94), "Field intelligence", "Attach records, incidents, inspections, and photos to locations.", WHITE, title_color=INK, body_color=MUTED)
    add_card(slide, Inches(9.82), Inches(3.1), Inches(2.1), Inches(1.94), "Spatial history", "Review what happened, where, and how often.", FOG, title_color=INK, body_color=MUTED)


def workflow(slide, num):
    add_page(slide, num, fill=FOG)
    add_label(slide, Inches(1.0), Inches(1.02), "Workflow", fill=GREEN, color=INK)
    add_text(slide, Inches(1.0), Inches(1.56), Inches(6.2), Inches(0.82), "The operating rhythm is map, route, execute, and review.", size=28, color=INK, bold=True)
    add_rect(slide, Inches(1.0), Inches(4.1), Inches(11.0), Inches(0.08), INK)
    steps = [
        (1.12, "01", "Map", "Define assets, boundaries, service zones, and active layers.", CYAN),
        (4.02, "02", "Route", "Build travel plans, dispatch logic, and coverage assignments.", GREEN),
        (6.92, "03", "Execute", "Track field work, deliveries, inspections, and exceptions.", AMBER),
        (9.82, "04", "Review", "Use spatial history to improve speed, coverage, and cost.", ROSE),
    ]
    for x, no, title, body, color in steps:
        add_rect(slide, Inches(x), Inches(3.86), Inches(0.18), Inches(0.54), color)
        add_text(slide, Inches(x + 0.28), Inches(3.28), Inches(2.0), Inches(0.22), no, size=11, color=color, bold=True)
        add_text(slide, Inches(x + 0.28), Inches(4.54), Inches(2.0), Inches(0.28), title, size=20, bold=True)
        add_text(slide, Inches(x + 0.28), Inches(4.98), Inches(2.02), Inches(0.86), body, size=12, color=MUTED)


def records(slide, num):
    add_page(slide, num, fill=WHITE)
    add_label(slide, Inches(1.0), Inches(1.02), "Records", fill=ROSE, color=WHITE)
    add_text(slide, Inches(1.0), Inches(1.56), Inches(6.5), Inches(0.82), "Every movement and field action should create usable spatial history.", size=28, color=INK, bold=True)
    add_text(slide, Inches(1.0), Inches(2.46), Inches(6.4), Inches(0.96), "Algomap preserves route history, field incidents, inspections, asset changes, and operator actions so teams can compare real geography-linked outcomes.", size=15, color=MUTED)
    add_rect(slide, Inches(7.78), Inches(0.96), Inches(4.2), Inches(5.62), PANEL)
    add_card(slide, Inches(8.08), Inches(1.3), Inches(3.6), Inches(1.02), "Route history", "Stops, timing, and movement stay tied to geography.", WHITE, title_color=INK, body_color=MUTED)
    add_card(slide, Inches(8.08), Inches(2.78), Inches(3.6), Inches(1.02), "Field history", "Inspections, incidents, and maintenance remain linked to locations.", CYAN, title_color=INK, body_color=INK)
    add_card(slide, Inches(8.08), Inches(4.26), Inches(3.6), Inches(1.02), "Decision history", "Teams can review what happened, where, and how the system responded.", GREEN, title_color=INK, body_color=INK)


def customers(slide, num):
    add_page(slide, num, fill=MIST)
    add_label(slide, Inches(1.0), Inches(1.02), "Customers", fill=SLATE, color=WHITE)
    add_text(slide, Inches(1.0), Inches(1.56), Inches(6.9), Inches(0.82), "Algomap serves any operation where movement, assets, and geography define performance.", size=28, color=INK, bold=True)
    add_card(slide, Inches(1.0), Inches(3.18), Inches(2.78), Inches(1.92), "Logistics teams", "Route vehicles, stops, schedules, and service coverage.", WHITE, title_color=INK, body_color=MUTED)
    add_card(slide, Inches(4.0), Inches(3.18), Inches(2.78), Inches(1.92), "Field services", "Coordinate inspections, repairs, and mobile crews spatially.", FOG, title_color=INK, body_color=MUTED)
    add_card(slide, Inches(7.0), Inches(3.18), Inches(2.78), Inches(1.92), "Infrastructure operators", "Monitor assets, zones, outages, and response patterns.", WHITE, title_color=INK, body_color=MUTED)
    add_card(slide, Inches(10.0), Inches(3.18), Inches(2.1), Inches(1.92), "Public programs", "Track service delivery, coverage, and reporting by area.", FOG, title_color=INK, body_color=MUTED)


def moat(slide, num):
    add_page(slide, num, fill=WHITE)
    add_label(slide, Inches(1.0), Inches(1.02), "Differentiation", fill=GREEN, color=INK)
    add_text(slide, Inches(1.0), Inches(1.56), Inches(6.7), Inches(0.82), "The moat is operational coordination around geography, not just maps on a screen.", size=28, color=INK, bold=True)
    add_rect(slide, Inches(1.0), Inches(2.72), Inches(11.08), Inches(0.06), INK)
    add_card(slide, Inches(1.0), Inches(3.22), Inches(3.46), Inches(2.02), "Spatial workflow", "Maps, routing, field execution, and records live together.", WHITE, title_color=INK, body_color=MUTED)
    add_card(slide, Inches(4.82), Inches(3.22), Inches(3.46), Inches(2.02), "Operational timing", "Dispatch, incidents, and changes stay close to the map context.", MIST, title_color=INK, body_color=MUTED)
    add_card(slide, Inches(8.64), Inches(3.22), Inches(3.46), Inches(2.02), "Persistent spatial memory", "History accumulates by zone, route, asset, and operator.", WHITE, title_color=INK, body_color=MUTED)


def business(slide, num):
    add_page(slide, num, fill=NIGHT)
    add_rect(slide, Inches(0.72), Inches(0.72), Inches(11.9), Inches(0.02), CYAN)
    add_rect(slide, Inches(0.72), Inches(6.82), Inches(11.9), Inches(0.02), CYAN)
    add_label(slide, Inches(1.0), Inches(1.02), "Business model", fill=WHITE, color=INK)
    add_text(slide, Inches(1.0), Inches(1.56), Inches(6.2), Inches(0.82), "Software revenue, operational deployment, and location intelligence tooling.", size=28, color=WHITE, bold=True)
    add_card(slide, Inches(1.0), Inches(3.26), Inches(3.48), Inches(1.96), "Platform subscriptions", "Charge by team, fleet, region, or active operational environment.", CYAN, title_color=INK, body_color=INK)
    add_card(slide, Inches(4.82), Inches(3.26), Inches(3.48), Inches(1.96), "Enterprise deployments", "Support larger logistics, utility, and infrastructure operators.", GREEN, title_color=INK, body_color=INK)
    add_card(slide, Inches(8.64), Inches(3.26), Inches(3.48), Inches(1.96), "Implementation and support", "Rollout, spatial modeling, and workflow design for complex operations.", AMBER, title_color=INK, body_color=INK)


def roadmap(slide, num):
    add_page(slide, num, fill=FOG)
    add_label(slide, Inches(1.0), Inches(1.02), "Roadmap", fill=AMBER, color=INK)
    add_text(slide, Inches(1.0), Inches(1.56), Inches(6.6), Inches(0.82), "The build path starts with spatial operations and expands into deeper optimization.", size=28, color=INK, bold=True)
    add_rect(slide, Inches(1.0), Inches(4.02), Inches(11.06), Inches(0.08), INK)
    phases = [
        (1.12, "Phase 1", "Core mapping, routing, assets, and field record modules.", CYAN),
        (4.72, "Phase 2", "Dispatch intelligence, exception workflows, and stronger analytics.", GREEN),
        (8.32, "Phase 3", "Predictive routing, infrastructure planning, and richer spatial automation.", AMBER),
    ]
    for x, title, body, color in phases:
        add_rect(slide, Inches(x), Inches(3.8), Inches(0.16), Inches(0.56), color)
        add_text(slide, Inches(x + 0.26), Inches(2.96), Inches(2.8), Inches(0.28), title, size=20, bold=True)
        add_text(slide, Inches(x + 0.26), Inches(4.44), Inches(2.64), Inches(0.86), body, size=12, color=MUTED)


def closing(slide, num):
    add_page(slide, num, fill=WHITE)
    add_rect(slide, Inches(0.88), Inches(0.96), Inches(11.56), Inches(5.92), PANEL)
    add_logo(slide, Inches(1.2), Inches(1.34), 0.82)
    add_label(slide, Inches(1.2), Inches(2.34), "Closing", fill=WHITE, color=INK)
    add_text(slide, Inches(1.2), Inches(2.92), Inches(6.36), Inches(0.82), "Algomap gives spatial operations one cleaner decision layer.", size=31, color=WHITE, bold=True)
    add_text(slide, Inches(1.2), Inches(3.9), Inches(6.3), Inches(1.0), "The platform connects GIS, routing, field records, and logistics performance so teams can move faster and plan with more confidence.", size=18, color=MIST)
    add_rect(slide, Inches(8.28), Inches(1.4), Inches(3.42), Inches(1.0), CYAN)
    add_rect(slide, Inches(8.86), Inches(2.72), Inches(2.84), Inches(0.82), GREEN)
    add_rect(slide, Inches(8.28), Inches(3.94), Inches(3.42), Inches(1.8), WHITE)
    add_text(slide, Inches(8.54), Inches(4.36), Inches(2.8), Inches(0.32), "Algomap", size=22, color=INK, bold=True)
    add_text(slide, Inches(8.54), Inches(4.86), Inches(2.76), Inches(0.66), "GIS, mapping, and logistics operating system for fleets, assets, and field teams.", size=12, color=MUTED)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]

    slides = [
        cover,
        problem,
        why_now,
        product,
        components,
        logistics,
        gis,
        workflow,
        records,
        customers,
        moat,
        business,
        roadmap,
        closing,
    ]

    for idx, fn in enumerate(slides, start=1):
        slide = prs.slides.add_slide(layout)
        fn(slide, idx)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
